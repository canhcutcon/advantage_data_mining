# -*- coding: utf-8 -*-
"""Generate editable PPTX mirroring main.tex (NeurIPS colors, 16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x05, 0x96, 0x69)

SW, SH = Inches(13.33), Inches(7.5)

# (title, bullets, image, note)
# bullet: (text, level, color_or_None, bold)
SLIDES = [
    ("TITLE", None, None,
     "Chờ giới thiệu. 'Em xin chào thầy/cô và các bạn. Hôm nay em trình bày đề tài "
     "Khai thác các mẫu hiếm — Rare Pattern Mining.' (0:15)"),

    ("Nội dung trình bày", [
        ("Bài toán & động cơ", 0, None, False),
        ("Ba định nghĩa mẫu hiếm", 0, None, False),
        ("Ba thuật toán: AprioriRare, AprioriInverse, CORI", 0, None, True),
        ("Độ đo tương quan & tính null-invariance", 0, None, False),
        ("Cài đặt, kiểm chứng tự động, thực nghiệm UCI Mushroom", 0, None, False),
    ], None,
     "Lướt nhanh 5 phần. Chuyển: 'Bắt đầu từ câu hỏi vì sao cần mẫu hiếm.' (0:20)"),

    ("Động cơ: mẫu hiếm nhưng giá trị cao", [
        ("Y tế: tổ hợp triệu chứng bệnh hiếm → giá trị chẩn đoán cao", 0, None, False),
        ("Công nghiệp: lỗi sản phẩm hiếm gặp = đối tượng cần phát hiện", 0, None, False),
        ("An ninh mạng: bất thường = hiếm theo định nghĩa", 0, None, False),
        ("FIM truyền thống bỏ qua hoàn toàn các mẫu này", 0, RED, True),
    ], None,
     "Mở bằng câu hỏi: 'Điều thú vị nhất trong dữ liệu có phải là điều xảy ra "
     "thường xuyên nhất không?' (1:00)"),

    ("Vấn đề: hạ minsup là không đủ", [
        ("Giải pháp ngây thơ: hạ thấp minsup", 0, None, False),
        ("Hệ quả 1: bùng nổ tổ hợp số itemset", 0, None, True),
        ("Hệ quả 2: nhiễu — nhiều 'mẫu hiếm' có sup = 0 (không tồn tại!)", 0, None, False),
        ("Cần: định nghĩa chặt chẽ + thuật toán riêng + độ đo lọc mẫu giả", 0, ACCENT, False),
    ], None,
     "Hạ minsup làm dàn itemset bùng nổ; nhiều itemset không tồn tại trong dữ liệu. "
     "Chuyển: 'Để nói chuyện cụ thể, em dùng một ví dụ nhỏ.' (1:00)"),

    ("Ví dụ chạy xuyên suốt", [
        ("T1: pasta, lemon, bread, orange", 0, None, False),
        ("T2: pasta, lemon", 0, None, False),
        ("T3: pasta, orange, cake", 0, None, False),
        ("T4: pasta, lemon, orange, cake", 0, None, False),
        ("minsup = 2 → 11 itemset phổ biến", 0, ACCENT, True),
        ("{bread}: sup = 1   •   {bread, cake}: sup = 0", 0, None, False),
    ], None,
     "CSDL 4 giao dịch theo bài giảng. Pasta có mặt mọi giao dịch — chi tiết này "
     "quay lại ở phần mẫu giả. (0:45)"),

    ("Ba định nghĩa mẫu hiếm", [
        ("Infrequent: sup < minsup — bùng nổ, chứa cả sup = 0", 0, None, False),
        ("Minimal rare (mRI): hiếm, mọi tập con đều phổ biến", 0, None, True),
        ("→ biên dưới của vùng hiếm trong dàn itemset", 1, ACCENT, False),
        ("Perfectly rare (PRI): minsup ≤ sup < maxsup, mọi item thành phần đều hiếm", 0, None, True),
        ("Ví dụ (minsup=2): mRI = {bread}, {lemon, cake}   •   PRI = {bread}", 0, None, False),
    ], None,
     "Tâm điểm lý thuyết — hình ảnh 'đường biên' trong dàn itemset. mRI nhỏ gọn mà "
     "đại diện mọi mẫu hiếm; PRI thêm ngưỡng dưới loại nhiễu. (1:15)"),

    ("AprioriRare: tìm biên hiếm tối thiểu", [
        ("Duyệt theo mức (level-wise) như Apriori", 0, None, False),
        ("Ứng viên sống sót qua prune mà không phổ biến → chính là mRI", 0, None, True),
        ("Itemset hiếm không dùng sinh ứng viên tiếp", 0, None, False),
        ("Ví dụ: mRI = {bread}, {lemon, cake}  ✓ khớp bài giảng", 0, GREEN, False),
        ("Nhược: phải duyệt toàn bộ vùng phổ biến", 0, RED, True),
    ], None,
     "Qua prune nghĩa là mọi tập con đều phổ biến → không phổ biến là tự động thành "
     "mRI. Nhược điểm chi phí định lượng ở Exp A. (1:15)"),

    ("AprioriInverse: thu hẹp ngay từ đầu", [
        ("Hai ngưỡng (minsup, maxsup)", 0, None, False),
        ("Loại mọi item phổ biến ngay bước khởi tạo", 0, None, True),
        ("Đúng đắn: mọi item hiếm → mọi tập con hiếm (điều kiện PRI tự thoả)", 0, None, False),
        ("Đổi lại: bỏ sót mẫu lai {item hiếm + item phổ biến}", 0, RED, False),
        ("Ví dụ (1.1, 3.1): 5 PRI  ✓ khớp bài giảng", 0, GREEN, False),
    ], None,
     "Không gian nhỏ hơn hẳn; cái giá là bỏ sót mẫu lai như "
     "{bệnh hiếm + triệu chứng phổ biến}. (1:00)"),

    ("Mẫu giả: phổ biến ≠ đáng tin", [
        ("{pasta, cake} xuất hiện ở 50% giao dịch", 0, None, True),
        ("Nhưng pasta có mặt trong mọi giao dịch", 0, None, False),
        ("Đồng xuất hiện là tất yếu — không có liên hệ thật (spurious)", 0, RED, False),
        ("Cần độ đo tương quan để lọc", 0, ACCENT, False),
    ], None,
     "Chuyển mạch sang nửa sau: hiếm thôi chưa đủ, còn cần 'thật'. (0:45)"),

    ("Hai độ đo: bond và all-confidence", [
        ("bond(X) = sup(X) / dsup(X) — tỉ lệ giao dịch 'liên quan' chứa toàn bộ X", 0, None, False),
        ("bond = 1: các mục luôn đi cùng nhau", 0, None, True),
        ("allconf(X) = sup(X) / max sup(item) — confidence nhỏ nhất của mọi luật từ X", 0, None, False),
        ("Cả hai: = 1 với itemset đơn; anti-monotone → cắt tỉa kiểu Apriori", 0, ACCENT, True),
        ("Ví dụ: bond({cake, bread}) = 0   •   bond({pasta, orange}) = 3/4", 0, None, False),
    ], None,
     "Nhấn tính anti-monotone — CORI dùng nó để cắt tỉa cả nhánh. (1:00)"),

    ("Null-invariance: vì sao lift thất bại", [
        ("Null transaction: giao dịch không chứa mục nào của mẫu", 0, None, False),
        ("χ², lift bị bóp méo bởi hàng triệu giao dịch không liên quan", 0, RED, False),
        ("bond, all-confidence, cosine, Kulczynski: null-invariant", 0, None, True),
        ("Mẫu hiếm ⊂ phần rất nhỏ dữ liệu → bắt buộc dùng null-invariant", 0, ACCENT, True),
    ], None,
     "Mẫu hiếm chỉ chạm phần rất nhỏ dữ liệu, còn lại toàn null transactions — "
     "độ đo nhạy với chúng bị bóp méo trên dữ liệu thưa. (0:45)"),

    ("CORI: hiếm + tương quan trong một lần duyệt", [
        ("Rare correlated: sup(X) < maxsup VÀ bond(X) ≥ minbond", 0, None, True),
        ("Hai ràng buộc đối ngẫu: hiếm = monotone • bond = anti-monotone", 0, None, False),
        ("Duyệt kiểu Eclat, mỗi itemset giữ TID-List + DTID-List", 0, None, False),
        ("bond rớt ngưỡng → cắt cả nhánh; sup ≥ maxsup → chưa xuất nhưng vẫn mở rộng", 0, None, False),
        ("Ví dụ (3, 0.6): {bread}, {cake}, {orange, cake}  ✓ khớp bài giảng", 0, GREEN, False),
    ], None,
     "Khi nối hai itemset: TID-List lấy giao, DTID-List lấy hợp — support và bond "
     "tính ngay. Nút chưa hiếm vẫn phải mở rộng. (1:15)"),

    ("Cài đặt & kiểm chứng 100%", [
        ("Python 3.13, chỉ thư viện chuẩn; TID-set dùng chung cho cả 3 thuật toán", 0, None, False),
        ("19 assert đối chiếu từng con số bài giảng — 100% ĐẠT", 0, GREEN, True),
        ("Phát hiện thiếu sót trong slide gốc: FIM liệt kê 10, đúng là 11", 0, None, True),
        ("(sót {pasta, orange, cake}, sup = 2)", 1, None, False),
        ("{pasta, cake}: lift = 1.0, bond = 0.5 → xác nhận mẫu giả bằng số", 0, None, False),
    ], None,
     "Kiểm chứng tự động đủ chặt để bắt được lỗi của chính tài liệu gốc. (1:00)"),

    ("Exp A — Cái giá của việc đi xuyên vùng phổ biến", [
        ("UCI Mushroom: 8 124 giao dịch, 118 item (dữ liệu dày)", 0, None, False),
        ("minsup 0.6 → 0.2: #frequent tăng ~890× (51 → 45 391); #mRI chỉ ~8× (122 → 986)", 0, None, True),
    ], "figures/mushroom_runtime.png",
     "Thời gian AprioriRare tăng theo cấp mũ khi minsup giảm. mRI gọn nhưng chi phí "
     "bị chi phối bởi vùng phổ biến. (1:00)"),

    ("Exp B & C — PRI rẻ, CORI giàu ngữ nghĩa", [
        ("AprioriInverse: mili-giây — alphabet hiếm rất nhỏ (không phải benchmark cùng điều kiện)", 0, None, False),
        ("CORI: bão hòa 90–104 mẫu khi nới minbond 0.95 → 0.5; mẫu bond = 1.0: {odor=m, …}, sup = 36 (0.44%)", 0, None, True),
    ], "figures/mushroom_experiments.png",
     "Mẫu đắt giá nhất: 36 cây nấm mùi mốc luôn đồng thời không vòng cuống + cuống "
     "màu quế — bond = 1 tuyệt đối, FIM không bao giờ chạm tới. (1:15)"),

    ("Kết luận: ba đánh đổi", [
        ("mRI gọn nhưng đắt — phải duyệt toàn vùng phổ biến", 0, None, True),
        ("PRI rẻ nhưng hẹp — bỏ sót mẫu lai", 0, None, True),
        ("CORI cân bằng — ràng buộc kép, kết quả nhỏ & giàu ngữ nghĩa", 0, None, True),
        ("Hướng 2021–2026: metaheuristic (MRI-CE), fuzzy, utility, privacy", 0, ACCENT, False),
    ], None,
     "Thông điệp: không có định nghĩa hiếm 'đúng' duy nhất — chỉ có đánh đổi; "
     "null-invariance là bắt buộc trên dữ liệu thưa. (0:50)"),

    ("THANKYOU", None, None,
     "Cảm ơn, mời câu hỏi. Chuẩn bị sẵn Q&A trong TALK_SCRIPT.md. (0:20)"),
]


def add_note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    for title, bullets, image, note in SLIDES:
        slide = prs.slides.add_slide(blank)

        if title == "TITLE":
            tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(2.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Khai thác các mẫu hiếm (Rare Pattern Mining)"
            p.alignment = PP_ALIGN.CENTER
            p.font.size, p.font.bold, p.font.color.rgb = Pt(40), True, PRIMARY
            p2 = tf.add_paragraph()
            p2.text = "Bài tập lớn môn Khai thác dữ liệu nâng cao"
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size, p2.font.color.rgb = Pt(24), ACCENT
            p3 = tf.add_paragraph()
            p3.text = "[Họ tên sinh viên — MSSV] · Tháng 6 năm 2026"
            p3.alignment = PP_ALIGN.CENTER
            p3.font.size, p3.font.color.rgb = Pt(18), DARK
            add_note(slide, note)
            continue

        if title == "THANKYOU":
            tb = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.33), Inches(3.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Xin cảm ơn thầy/cô và các bạn!"
            p.alignment = PP_ALIGN.CENTER
            p.font.size, p.font.bold, p.font.color.rgb = Pt(36), True, PRIMARY
            for line, size in [
                ("3 định nghĩa • 3 thuật toán • kiểm chứng 100% • UCI Mushroom", 20),
                ("Mã nguồn & kết quả: advantage_data_mining/code, results/", 16),
                ("Mời thầy/cô và các bạn đặt câu hỏi", 26),
            ]:
                q = tf.add_paragraph()
                q.text = line
                q.alignment = PP_ALIGN.CENTER
                q.font.size, q.font.color.rgb = Pt(size), DARK
            add_note(slide, note)
            continue

        # Frame title
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size, p.font.bold, p.font.color.rgb = Pt(30), True, PRIMARY

        body_top = Inches(1.4)
        if image:
            pic_h = Inches(4.3)
            slide.shapes.add_picture(image, Inches(1.8), body_top, height=pic_h)
            body_top = Inches(5.9)
            body_h = Inches(1.4)
        else:
            body_h = Inches(5.6)

        tb = slide.shapes.add_textbox(Inches(0.8), body_top, Inches(11.7), body_h)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for text, level, color, bold in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = ("• " if level == 0 else "   – ") + text
            p.font.size = Pt(15 if image else 20)
            p.font.bold = bold
            p.font.color.rgb = color if color else DARK
            p.space_after = Pt(6 if image else 12)
        add_note(slide, note)

    prs.save("presentation.pptx")
    print(f"OK: presentation.pptx — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
