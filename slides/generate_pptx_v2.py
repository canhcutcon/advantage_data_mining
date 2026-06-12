# -*- coding: utf-8 -*-
"""Black & Gold redesign of the Rare Pattern Mining deck (same 17-slide content).

Reuses SLIDES content from generate_pptx.py; remaps the NeurIPS palette to:
gold #BF9A4A, charcoal #191919, cream #FAF7F0, forest #40695B, rust #C15937.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import generate_pptx as v1
from generate_pptx import SLIDES

GOLD = RGBColor(0xBF, 0x9A, 0x4A)
GOLD_DARK = RGBColor(0x99, 0x79, 0x29)
CHARCOAL = RGBColor(0x19, 0x19, 0x19)
CREAM = RGBColor(0xFA, 0xF7, 0xF0)
CREAM_TEXT = RGBColor(0xF4, 0xF6, 0xF6)
BODY = RGBColor(0x2C, 0x2C, 0x2C)
FOREST = RGBColor(0x40, 0x69, 0x5B)
RUST = RGBColor(0xC1, 0x59, 0x37)

# old palette -> new palette
REMAP = {v1.PRIMARY: GOLD_DARK, v1.ACCENT: GOLD_DARK, v1.RED: RUST,
         v1.GREEN: FOREST, v1.DARK: BODY}

SW, SH = Inches(13.33), Inches(7.5)
BAND_W = Inches(0.45)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_scaffold(slide, idx_label):
    """Charcoal side band + slide number in gold."""
    rect(slide, 0, 0, BAND_W, SH, CHARCOAL)
    tb = slide.shapes.add_textbox(0, Inches(6.9), BAND_W, Inches(0.45))
    tb.text_frame.margin_left = tb.text_frame.margin_right = 0
    p = tb.text_frame.paragraphs[0]
    p.text = idx_label
    p.alignment = PP_ALIGN.CENTER
    p.font.size, p.font.bold, p.font.color.rgb = Pt(11), True, GOLD


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]
    total = len(SLIDES)

    for i, (title, bullets, image, note) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)

        if title == "TITLE":
            bg(slide, CHARCOAL)
            rect(slide, Inches(2.665), Inches(2.05), Inches(8.0), Pt(3), GOLD)
            tb = slide.shapes.add_textbox(Inches(1), Inches(2.35), Inches(11.33), Inches(3.0))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Khai thác các mẫu hiếm"
            p.alignment = PP_ALIGN.CENTER
            p.font.size, p.font.bold, p.font.color.rgb = Pt(44), True, GOLD
            p2 = tf.add_paragraph()
            p2.text = "(Rare Pattern Mining)"
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size, p2.font.bold, p2.font.color.rgb = Pt(28), True, GOLD
            p3 = tf.add_paragraph()
            p3.text = "Bài tập lớn môn Khai thác dữ liệu nâng cao"
            p3.alignment = PP_ALIGN.CENTER
            p3.font.size, p3.font.color.rgb = Pt(20), CREAM_TEXT
            p3.space_before = Pt(18)
            p4 = tf.add_paragraph()
            p4.text = "[Họ tên sinh viên — MSSV] · Tháng 6 năm 2026"
            p4.alignment = PP_ALIGN.CENTER
            p4.font.size, p4.font.color.rgb = Pt(15), GOLD
            p4.space_before = Pt(10)
            rect(slide, Inches(2.665), Inches(5.45), Inches(8.0), Pt(3), GOLD)
            add_note(slide, note)
            continue

        if title == "THANKYOU":
            bg(slide, CHARCOAL)
            rect(slide, Inches(2.665), Inches(1.7), Inches(8.0), Pt(3), GOLD)
            tb = slide.shapes.add_textbox(Inches(1), Inches(2.1), Inches(11.33), Inches(3.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Xin cảm ơn thầy/cô và các bạn!"
            p.alignment = PP_ALIGN.CENTER
            p.font.size, p.font.bold, p.font.color.rgb = Pt(38), True, GOLD
            for line, size, col, before in [
                ("3 định nghĩa ◆ 3 thuật toán ◆ kiểm chứng 100% ◆ UCI Mushroom", 18, CREAM_TEXT, 24),
                ("Mã nguồn & kết quả: advantage_data_mining/code, results/", 14, CREAM_TEXT, 12),
                ("Mời thầy/cô và các bạn đặt câu hỏi", 26, GOLD, 30),
            ]:
                q = tf.add_paragraph()
                q.text = line
                q.alignment = PP_ALIGN.CENTER
                q.font.size, q.font.color.rgb = Pt(size), col
                q.space_before = Pt(before)
            rect(slide, Inches(2.665), Inches(5.8), Inches(8.0), Pt(3), GOLD)
            add_note(slide, note)
            continue

        # ----- content slide -----
        bg(slide, CREAM)
        content_scaffold(slide, f"{i}/{total}")

        # title + gold rule
        tb = slide.shapes.add_textbox(Inches(0.85), Inches(0.35), Inches(12.0), Inches(0.85))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size, p.font.bold, p.font.color.rgb = Pt(29), True, CHARCOAL
        rect(slide, Inches(0.9), Inches(1.18), Inches(4.2), Pt(3.5), GOLD)

        body_top = Inches(1.5)
        if image:
            # fit image inside a max box, centered over the content area
            from PIL import Image as PILImage
            iw, ih = PILImage.open(image).size
            max_w, max_h = 11.3, 4.2
            w_in = min(max_w, iw * max_h / ih)
            h_in = w_in * ih / iw
            left = Inches(0.45 + (12.88 - w_in) / 2)
            slide.shapes.add_picture(image, left, body_top, width=Inches(w_in))
            body_top = Inches(1.55 + h_in + 0.15)
            body_h = Inches(7.3 - (1.55 + h_in + 0.15))
        else:
            body_h = Inches(5.5)

        tb = slide.shapes.add_textbox(Inches(1.0), body_top, Inches(11.6), body_h)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for text, level, color, bold in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            marker = "◆  " if level == 0 else "      –  "
            r1 = p.add_run()
            r1.text = marker
            r1.font.size = Pt(15 if image else 20)
            r1.font.color.rgb = GOLD_DARK if level == 0 else BODY
            r1.font.bold = level == 0
            r2 = p.add_run()
            r2.text = text
            r2.font.size = Pt(15 if image else 20)
            r2.font.bold = bold
            r2.font.color.rgb = REMAP.get(color, BODY) if color else BODY
            p.space_after = Pt(6 if image else 12)
        add_note(slide, note)

    prs.save("presentation_v2.pptx")
    print(f"OK: presentation_v2.pptx — {total} slides (Black & Gold)")


if __name__ == "__main__":
    main()
